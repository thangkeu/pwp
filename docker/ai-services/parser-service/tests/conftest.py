"""
conftest.py — sinh fixture file NGAY TRONG LÚC TEST bằng chính các thư viện parse tương ứng,
thay vì commit file nhị phân mẫu vào repo (tránh phình repo, tránh fixture "giả" không phản ánh
đúng output thật của Word/Excel/PowerPoint/PDF/PIL).
"""

from __future__ import annotations

import base64
import io
import shutil
import urllib.parse
import zlib

import docx
import ezdxf
import pytest
import vsdx
from fpdf import FPDF
from openpyxl import Workbook
from PIL import Image, ImageDraw
from pptx import Presentation


@pytest.fixture
def sample_docx_bytes() -> bytes:
    document = docx.Document()
    document.add_heading("Báo giá Fortigate 100F", level=1)
    document.add_paragraph("Đây là tài liệu mẫu dùng cho unit test DocxAdapter.")
    table = document.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Part Number"
    table.rows[0].cells[1].text = "Qty"
    table.rows[1].cells[0].text = "FG-100F"
    table.rows[1].cells[1].text = "2"

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


@pytest.fixture
def sample_xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "BOM"
    sheet.append(["Part Number", "Qty", "Unit Price"])
    sheet.append(["FG-100F", 2, 15000000])

    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


@pytest.fixture
def sample_pptx_bytes() -> bytes:
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Kiến trúc SD-WAN đề xuất"
    body = slide.placeholders[1]
    body.text_frame.text = "Nội dung mẫu cho unit test PptxAdapter"

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


@pytest.fixture
def sample_pdf_bytes() -> bytes:
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(0, 10, text="Tai lieu mau cho unit test PdfAdapter")
    return bytes(pdf.output())


@pytest.fixture
def sample_png_bytes() -> bytes:
    image = Image.new("RGB", (300, 80), color="white")
    draw = ImageDraw.Draw(image)
    draw.text((10, 30), "HELLO OCR TEST", fill="black")
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


@pytest.fixture
def sample_markdown_bytes() -> bytes:
    content = (
        "# Tiêu đề\n\n"
        "Đoạn văn có [link](https://example.com/docs).\n\n"
        "| A | B |\n|---|---|\n| 1 | 2 |\n"
    )
    return content.encode("utf-8")


@pytest.fixture
def sample_html_bytes() -> bytes:
    content = (
        "<html><body><h1>Tiêu đề</h1>"
        "<p>Xem <a href='https://example.com'>tại đây</a></p>"
        "<table><tr><td>1</td><td>2</td></tr></table>"
        "</body></html>"
    )
    return content.encode("utf-8")


@pytest.fixture
def sample_dxf_bytes() -> bytes:
    doc = ezdxf.new("R2010")
    msp = doc.modelspace()
    msp.add_text("Fortigate 100F - Rack A1", dxfattribs={"layer": "TEXT"}).set_placement((0, 0))
    msp.add_mtext("Ghi chu: 2x FG-100F HA cluster", dxfattribs={"layer": "NOTES"}).set_location((0, 10))

    buffer = io.StringIO()
    doc.write(buffer, fmt="asc")
    return buffer.getvalue().encode("utf-8")


@pytest.fixture
def sample_vsdx_bytes(tmp_path) -> bytes:
    # Thư viện vsdx bắt buộc thao tác trên file thật (không nhận BytesIO) — dùng template mẫu
    # đi kèm chính thư viện `vsdx` (media/media.vsdx) rồi sửa text thành nội dung mẫu presales,
    # thay vì commit 1 file .vsdx nhị phân riêng vào repo test.
    template_path = vsdx.__file__.replace("__init__.py", "media/media.vsdx")
    tmp_file = tmp_path / "sample.vsdx"
    shutil.copy(template_path, tmp_file)

    with vsdx.VisioFile(str(tmp_file)) as visio_file:
        page = visio_file.pages[0]
        page.find_replace("RECTANGLE", "Kien truc SD-WAN de xuat")
        visio_file.save_vsdx(str(tmp_file))

    return tmp_file.read_bytes()


@pytest.fixture
def sample_drawio_bytes() -> bytes:
    """Dạng XML thô (không nén) — <diagram> chứa trực tiếp <mxGraphModel>."""
    content = """<mxfile host="app.diagrams.net">
  <diagram name="Kien truc SD-WAN" id="abc123">
    <mxGraphModel dx="800" dy="600">
      <root>
        <mxCell id="0"/>
        <mxCell id="1" parent="0"/>
        <mxCell id="2" value="Fortigate 100F" vertex="1" parent="1">
          <mxGeometry x="40" y="40" width="120" height="60" as="geometry"/>
        </mxCell>
        <mxCell id="3" value="&lt;b&gt;Branch Office&lt;/b&gt;&lt;br&gt;200 users" vertex="1" parent="1">
          <mxGeometry x="200" y="40" width="160" height="80" as="geometry"/>
        </mxCell>
        <mxCell id="4" value="ket noi VPN" edge="1" parent="1" source="2" target="3">
          <mxGeometry relative="1" as="geometry"/>
        </mxCell>
      </root>
    </mxGraphModel>
  </diagram>
</mxfile>"""
    return content.encode("utf-8")


@pytest.fixture
def sample_drawio_compressed_bytes() -> bytes:
    """Dạng nén base64(deflate(urlencode(xml))) — mặc định khi lưu từ app.diagrams.net."""
    inner_xml = (
        '<mxGraphModel><root><mxCell id="0"/><mxCell id="1" parent="0"/>'
        '<mxCell id="2" value="Node nen thu (compressed)" vertex="1" parent="1">'
        "<mxGeometry/></mxCell></root></mxGraphModel>"
    )
    encoded_uri = urllib.parse.quote(inner_xml)
    compressed = zlib.compress(encoded_uri.encode("utf-8"))[2:-4]  # raw deflate
    b64 = base64.b64encode(compressed).decode("ascii")
    content = f'<mxfile><diagram name="Nen" id="x">{b64}</diagram></mxfile>'
    return content.encode("utf-8")
