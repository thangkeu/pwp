"""PptxAdapter — trích xuất từ file .pptx dùng python-pptx."""

from __future__ import annotations

import base64
import io
import zipfile

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from app.adapters.base import ParserAdapter, ParsingError
from app.domain.parsed_document import (
    ParsedDocument,
    ParsedDocumentMetadata,
    ParsedImage,
)


class PptxAdapter(ParserAdapter):
    @property
    def supported_extensions(self) -> list[str]:
        return ["pptx"]

    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        try:
            presentation = Presentation(io.BytesIO(content))
        except (PackageNotFoundError, KeyError, ValueError, zipfile.BadZipFile) as exc:
            raise ParsingError(f"File PPTX không hợp lệ hoặc hỏng: {filename} ({exc})") from exc

        warnings: list[str] = []
        text_parts: list[str] = []
        images: list[ParsedImage] = []
        image_index = 0

        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide_texts.append(shape.text_frame.text.strip())

                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image = shape.image
                        images.append(
                            ParsedImage(
                                index=image_index,
                                page_number=slide_number,
                                content_base64=base64.b64encode(image.blob).decode("ascii"),
                            )
                        )
                        image_index += 1
                    except Exception as exc:  # noqa: BLE001
                        warnings.append(f"Slide {slide_number}: không trích được 1 hình ảnh ({exc})")

            notes = ""
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_texts.append(f"[Notes] {notes}")

            if slide_texts:
                text_parts.append(f"[Slide {slide_number}]\n" + "\n".join(slide_texts))

        text = "\n\n".join(text_parts)
        word_count = len(text.split())

        return ParsedDocument(
            text=text,
            images=images,
            metadata=ParsedDocumentMetadata(
                parser_name="PptxAdapter",
                source_filename=filename,
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                page_count=len(presentation.slides),
                word_count=word_count,
                warnings=warnings,
            ),
        )
